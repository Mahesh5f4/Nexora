import abc

class DocumentProcessor(abc.ABC):
    @abc.abstractmethod
    def can_process(self, content_type: str, filename: str) -> bool:
        pass

    @abc.abstractmethod
    def extract_text(self, content: bytes) -> str:
        pass

class PlainTextProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "text/plain" or filename.lower().endswith(".txt")

    def extract_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            # Fallback to ascii with replacement if not utf-8
            return content.decode("ascii", errors="replace")

class PdfProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "application/pdf" or filename.lower().endswith(".pdf")

    def extract_text(self, content: bytes) -> str:
        import io
        from pypdf import PdfReader
        
        pdf = PdfReader(io.BytesIO(content))
        text = ""
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text

class DocxProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.lower().endswith(".docx")

    def extract_text(self, content: bytes) -> str:
        import io
        import docx
        
        doc = docx.Document(io.BytesIO(content))
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)

class XlsxProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or filename.lower().endswith((".xlsx", ".xls"))

    def extract_text(self, content: bytes) -> str:
        import io
        import openpyxl
        
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text = []
        for sheet in wb.worksheets:
            text.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and join with tabs
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    text.append("\t".join(row_values))
        return "\n".join(text)

class CsvProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "text/csv" or filename.lower().endswith(".csv")

    def extract_text(self, content: bytes) -> str:
        import io
        import csv
        
        try:
            decoded = content.decode("utf-8")
        except UnicodeDecodeError:
            decoded = content.decode("ascii", errors="replace")
            
        reader = csv.reader(io.StringIO(decoded))
        text = []
        for row in reader:
            text.append("\t".join(row))
        return "\n".join(text)

class PptxProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or filename.lower().endswith(".pptx")

    def extract_text(self, content: bytes) -> str:
        import io
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(content))
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)

class MarkdownProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "text/markdown" or filename.lower().endswith(".md")

    def extract_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("ascii", errors="replace")

class HtmlProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "text/html" or filename.lower().endswith((".html", ".htm"))

    def extract_text(self, content: bytes) -> str:
        try:
            html = content.decode("utf-8")
        except UnicodeDecodeError:
            html = content.decode("ascii", errors="replace")
            
        from html.parser import HTMLParser
        
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.ignore = False
                
            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style'):
                    self.ignore = True
                    
            def handle_endtag(self, tag):
                if tag in ('script', 'style'):
                    self.ignore = False
                    
            def handle_data(self, data):
                if not self.ignore:
                    text = data.strip()
                    if text:
                        self.text.append(text)
                        
        extractor = TextExtractor()
        extractor.feed(html)
        return "\n".join(extractor.text)

class JsonProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type == "application/json" or filename.lower().endswith(".json")

    def extract_text(self, content: bytes) -> str:
        import json
        
        try:
            decoded = content.decode("utf-8")
            data = json.loads(decoded)
            return json.dumps(data, indent=2)
        except Exception:
            try:
                return content.decode("utf-8")
            except:
                return content.decode("ascii", errors="replace")

class XmlProcessor(DocumentProcessor):
    def can_process(self, content_type: str, filename: str) -> bool:
        return content_type in ("application/xml", "text/xml") or filename.lower().endswith(".xml")

    def extract_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("ascii", errors="replace")

